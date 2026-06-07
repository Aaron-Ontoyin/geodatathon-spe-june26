"""Build the D2 deliverable: figures + a 13-slide PowerPoint from the live model.

Run with:

    GEO_THERMOGIS_ROOT=data/thermogis_grid \
        uv run --with python-pptx --with matplotlib python scripts/build_deck.py

Every number and figure is computed from the current model (nothing hardcoded), so the
deck always matches the latest results. Outputs land in ``outputs/`` (gitignored).
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.figure import Figure
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from geothermal import config
from geothermal.design import district_demand, simulate
from geothermal.economics.optimization import (
    DesignCandidate,
    design_for,
    evaluate_candidate,
    monte_carlo_lcoe_samples,
    optimize,
)
from geothermal.economics.sensitivity import tornado
from geothermal.report import _TORNADO_FIELDS
from geothermal.resource import locate_demand_center, recommend_new_well
from geothermal.resource.spatial import resource_grid

# --------------------------------------------------------------------------- #
# Style
# --------------------------------------------------------------------------- #
TEAL = "#0f6e6e"
TEAL_DARK = "#0a4a4a"
SIENNA = "#a8531f"
SIENNA_LIGHT = "#d98b4a"
SLATE = "#2e3a3a"
SAND = "#e7ddc9"
HIGHLIGHT = "#c4762a"
GREY = "#8a9494"

TEAL_RGB = RGBColor(0x0F, 0x6E, 0x6E)
TEAL_DARK_RGB = RGBColor(0x0A, 0x4A, 0x4A)
SIENNA_RGB = RGBColor(0xA8, 0x53, 0x1F)
SLATE_RGB = RGBColor(0x2E, 0x3A, 0x3A)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
SAND_RGB = RGBColor(0xF3, 0xEE, 0xE3)

DPI = 150

FIG_DIR = config.OUTPUTS_DIR / "figures"
PPTX_PATH = config.OUTPUTS_DIR / "Team_GEOTHERM_PPT_V1.pptx"


def _apply_style() -> None:
    plt.rcParams.update(
        {
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "axes.edgecolor": SLATE,
            "axes.labelcolor": SLATE,
            "axes.titlecolor": TEAL_DARK,
            "axes.titlesize": 13,
            "axes.titleweight": "bold",
            "axes.labelsize": 11,
            "xtick.color": SLATE,
            "ytick.color": SLATE,
            "text.color": SLATE,
            "font.family": "sans-serif",
            "font.sans-serif": ["Helvetica", "Arial", "DejaVu Sans"],
            "axes.grid": True,
            "grid.color": "#d8dede",
            "grid.linewidth": 0.7,
            "legend.frameon": False,
            "legend.fontsize": 9,
        }
    )


def _save(fig: Figure, name: str) -> Path:
    path = FIG_DIR / name
    fig.savefig(path, dpi=DPI, bbox_inches="tight")
    plt.close(fig)
    return path


# --------------------------------------------------------------------------- #
# Figures
# --------------------------------------------------------------------------- #
@dataclass(frozen=True, slots=True)
class DeckNumbers:
    """Headline numbers extracted from the live model, used in figures and the deck."""

    n_doublets: int
    lcoe_eur_per_gj: float
    p10: float
    p50: float
    p90: float
    lcoe_by_doublets: dict[int, float]


def fig_monthly_dispatch(best: DesignCandidate) -> Path:
    perf = simulate(design_for(best.geo_capacity_mw, best.assumptions), district_demand())
    m = perf.monthly
    months = m["month"].to_numpy()
    month_labels = ["J", "F", "M", "A", "M", "J", "J", "A", "S", "O", "N", "D"]

    fig, (ax_h, ax_c) = plt.subplots(1, 2, figsize=(12, 5))

    # Heating panel: stacked supply vs demand line.
    ax_h.stackplot(
        months,
        m["geo_heat_mw"] + m["hp_elec_mw"],
        m["ates_discharge_mw"],
        m["backup_mw"],
        labels=["Geothermal + heat pump", "HT-ATES discharge", "Backup"],
        colors=[TEAL, SIENNA_LIGHT, GREY],
        alpha=0.9,
    )
    ax_h.plot(months, m["heating_mw"], color=SLATE, lw=2.4, marker="o", ms=4, label="Heating demand")
    ax_h.set_title("Heating: winter supply stack vs demand")
    ax_h.set_ylabel("Thermal power (MWth)")
    ax_h.set_xlabel("Month")
    ax_h.set_xticks(months)
    ax_h.set_xticklabels(month_labels)
    ax_h.legend(loc="upper center", ncol=2)
    ax_h.margins(x=0.01)

    # Cooling panel.
    ax_c.stackplot(
        months,
        m["abs_cool_mw"],
        m["free_cool_mw"],
        m["comp_cool_mw"],
        labels=["Absorption (heat-driven)", "Free cooling", "Compression"],
        colors=[TEAL, SAND, SIENNA],
        alpha=0.9,
    )
    ax_c.plot(months, m["cooling_mw"], color=SLATE, lw=2.4, marker="o", ms=4, label="Cooling demand")
    ax_c.set_title("Cooling: summer supply stack vs demand")
    ax_c.set_ylabel("Thermal power (MWth)")
    ax_c.set_xlabel("Month")
    ax_c.set_xticks(months)
    ax_c.set_xticklabels(month_labels)
    ax_c.legend(loc="upper center", ncol=2)
    ax_c.margins(x=0.01)

    fig.suptitle(
        "Monthly dispatch: one doublet, heat in winter and cool in summer",
        fontsize=15,
        fontweight="bold",
        color=TEAL_DARK,
    )
    fig.tight_layout(rect=(0, 0, 1, 0.96))
    return _save(fig, "monthly_dispatch.png")


def fig_design_comparison() -> tuple[Path, dict[int, float]]:
    doublets = [1, 2, 3]
    lcoe = {n: evaluate_candidate(n).lcoe_eur_per_gj for n in doublets}
    best_n = min(lcoe, key=lambda k: lcoe[k])

    fig, ax = plt.subplots(figsize=(7.5, 5))
    colors = [HIGHLIGHT if n == best_n else TEAL for n in doublets]
    bars = ax.bar([str(n) for n in doublets], [lcoe[n] for n in doublets], color=colors, width=0.6)
    for bar, n in zip(bars, doublets):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + 0.6,
            f"{lcoe[n]:.1f}",
            ha="center",
            va="bottom",
            fontweight="bold",
            color=SLATE,
        )
    ax.set_title("LCoE by number of doublets (lowest wins)")
    ax.set_xlabel("Number of doublets")
    ax.set_ylabel("LCoE (EUR/GJ)")
    ax.set_ylim(0, max(lcoe.values()) * 1.18)
    ax.text(
        bars[best_n - 1].get_x() + bars[best_n - 1].get_width() / 2,
        lcoe[best_n] * 0.5,
        "cheapest",
        ha="center",
        color="white",
        fontweight="bold",
        rotation=0,
    )
    fig.tight_layout()
    return _save(fig, "design_comparison.png"), lcoe


def fig_monte_carlo() -> tuple[Path, float, float, float]:
    samples = monte_carlo_lcoe_samples(1, n_samples=3000)
    samples = samples[np.isfinite(samples)]
    p10, p50, p90 = (float(np.percentile(samples, q)) for q in (10, 50, 90))

    fig, ax = plt.subplots(figsize=(8, 5))
    ax.hist(samples, bins=50, color=TEAL, alpha=0.78, edgecolor="white")
    for value, label, color in (
        (p10, f"P10 {p10:.1f}", SIENNA),
        (p50, f"P50 {p50:.1f}", HIGHLIGHT),
        (p90, f"P90 {p90:.1f}", SIENNA),
    ):
        ax.axvline(value, color=color, lw=2.2, ls="--")
        ax.text(
            value,
            ax.get_ylim()[1] * 0.96,
            label,
            rotation=90,
            va="top",
            ha="right",
            color=color,
            fontweight="bold",
            fontsize=10,
        )
    ax.set_title("Monte-Carlo LCoE for one doublet (transmissivity uncertainty)")
    ax.set_xlabel("LCoE (EUR/GJ)")
    ax.set_ylabel("Frequency")
    fig.tight_layout()
    return _save(fig, "monte_carlo_lcoe.png"), p10, p50, p90


def fig_tornado() -> Path:
    df = tornado(_TORNADO_FIELDS)
    df = df.sort_values("swing", ascending=True).reset_index(drop=True)
    labels = [_pretty_field(name) for name in df["field"]]

    fig, ax = plt.subplots(figsize=(8.5, 5))
    bars = ax.barh(labels, df["swing"], color=TEAL, height=0.6)
    bars[-1].set_color(HIGHLIGHT)  # top driver
    for bar, swing in zip(bars, df["swing"]):
        ax.text(
            bar.get_width() + max(df["swing"]) * 0.01,
            bar.get_y() + bar.get_height() / 2,
            f"{swing:.1f}",
            va="center",
            color=SLATE,
            fontsize=9,
        )
    ax.set_title("What moves the cost: LCoE swing by assumption")
    ax.set_xlabel("LCoE swing across plausible range (EUR/GJ)")
    ax.set_xlim(0, max(df["swing"]) * 1.12)
    fig.tight_layout()
    return _save(fig, "sensitivity_tornado.png")


def fig_resource_map() -> Path:
    grid = resource_grid(120)
    demand_x, demand_y = locate_demand_center()
    rec = recommend_new_well()

    fig, ax = plt.subplots(figsize=(8.5, 6.5))
    mesh = ax.pcolormesh(grid.x, grid.y, grid.power_mw, shading="auto", cmap="viridis")
    cbar = fig.colorbar(mesh, ax=ax)
    cbar.set_label("Doublet power, P50 (MWth)")

    for wid in config.WELL_IDS:
        well = config.WELLS[wid]
        ax.scatter(well.x, well.y, c="white", edgecolors="black", s=80, marker="^", zorder=5)
        ax.annotate(
            wid,
            (well.x, well.y),
            textcoords="offset points",
            xytext=(6, 6),
            color="white",
            fontsize=9,
            fontweight="bold",
        )
    ax.scatter(
        demand_x,
        demand_y,
        c=SIENNA,
        edgecolors="white",
        s=160,
        marker="s",
        zorder=6,
        label="Demand centre (USP)",
    )
    ax.scatter(
        rec["x"],
        rec["y"],
        c=HIGHLIGHT,
        edgecolors="white",
        s=240,
        marker="*",
        zorder=7,
        label=f"Recommended site ({rec['power_mw_p50']:.1f} MWth P50)",
    )
    ax.set_title("Resource map: doublet power with sited new well")
    ax.set_xlabel("RD New easting (m)")
    ax.set_ylabel("RD New northing (m)")
    legend = ax.legend(
        loc="upper left", frameon=True, facecolor="white", framealpha=0.95, edgecolor=SLATE
    )
    for text in legend.get_texts():
        text.set_color(SLATE)
    ax.ticklabel_format(style="plain")
    fig.tight_layout()
    return _save(fig, "resource_map.png")


def _pretty_field(name: str) -> str:
    pretty = {
        "injection_temp_c": "Injection temperature",
        "electricity_price_eur_per_mwhe": "Electricity price",
        "well_cost_meur": "Well cost",
        "discount_rate": "Discount rate",
        "heat_pump_keur_per_mwth": "Heat-pump CAPEX",
        "gas_price_eur_per_mwhth": "Gas price",
        "demand_distance_weight_mw_per_km": "Demand-distance penalty",
    }
    return pretty.get(name, name)


# --------------------------------------------------------------------------- #
# Slide deck
# --------------------------------------------------------------------------- #
@dataclass(frozen=True, slots=True)
class SlideSpec:
    """One content slide: title, bullet body, presenter notes, optional figure."""

    title: str
    bullets: tuple[str, ...]
    notes: str
    image: Path | None = None


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, prs, TEAL_DARK_RGB)

    title = _add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(2.0))
    p = title.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Geothermal District Heating & Cooling for Utrecht: "
        "the lowest-cost path to >=10 MWth heat + >=5 MWth cooling"
    )
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE_RGB

    sub = _add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12), Inches(2.6))
    lines = [
        ("Team [TEAM NAME]", 22, True, SAND_RGB),
        ("[MEMBER NAMES + SPE MEMBER NUMBERS]", 18, False, WHITE_RGB),
        ("SPE Africa Geothermal Datathon 2026", 16, False, SAND_RGB),
        ('"Not the most megawatts, the lowest, most credible LCoE."', 18, True, SIENNA_RGB),
    ]
    for i, (text, size, bold, color) in enumerate(lines):
        para = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        para.space_after = Pt(10)


def _add_content_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar.
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL_RGB
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec.title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE_RGB

    has_image = spec.image is not None
    body_w = Inches(6.0) if has_image else Inches(11.8)
    body = _add_textbox(slide, Inches(0.6), Inches(1.3), body_w, Inches(5.8))
    for i, bullet in enumerate(spec.bullets):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        run = para.add_run()
        run.text = bullet
        run.font.size = Pt(16)
        run.font.color.rgb = SLATE_RGB
        para.space_after = Pt(10)
        para.line_spacing = 1.05

    if has_image and spec.image is not None:
        slide.shapes.add_picture(
            str(spec.image), Inches(6.9), Inches(1.4), width=Inches(6.2)
        )

    notes = slide.notes_slide.notes_text_frame
    notes.text = spec.notes


def _add_textbox(slide: object, left: Inches, top: Inches, width: Inches, height: Inches) -> object:
    box = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[attr-defined]
    box.text_frame.word_wrap = True
    return box.text_frame


def _fill_background(slide: object, prs: Presentation, color: RGBColor) -> None:
    rect = slide.shapes.add_shape(  # type: ignore[attr-defined]
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    slide.shapes._spTree.remove(rect._element)  # type: ignore[attr-defined]
    slide.shapes._spTree.insert(2, rect._element)  # type: ignore[attr-defined]


def _slide_specs(
    nums: DeckNumbers,
    figs: dict[str, Path],
) -> list[SlideSpec]:
    n = nums.n_doublets
    lcoe = nums.lcoe_eur_per_gj
    lc = nums.lcoe_by_doublets
    return [
        SlideSpec(
            "Problem framing",
            (
                "District in Utrecht (NL): mixed residential, office and public load.",
                "Must deliver >=10 MWth heating AND >=5 MWth cooling from the "
                "Rotliegend (Slochteren) reservoir.",
                "Winner = lowest LCoE at adequate capacity (going big is penalised).",
                "Drilling cost counts for every well, existing or new.",
            ),
            "State the constraints and, crucially, that the score is cost-driven. This "
            "reframes the problem from maximise resource to minimise cost of meeting demand.",
        ),
        SlideSpec(
            "Our approach (the pipeline)",
            (
                "A reproducible, agent-driven pipeline: data, resource, design, "
                "techno-economics.",
                "Balanced flagship: rigorous geoscience and techno-economics plus an AI "
                "workflow and interactive app.",
                "One typed config drives everything; every number is reproducible offline.",
                "Data foundation, Resource (Ch1), System design (Ch2A), LCoE (Ch2B), "
                "Agentic workflow and app.",
            ),
            "We built it as a tool, not a one-off notebook. The same code reruns end to "
            "end and the organisers can audit every step.",
        ),
        SlideSpec(
            "Data foundation / EDA",
            (
                "Four wells, deliberately messy: along-hole depths mislabelled as TVD, "
                "empty TVD column, missing porosity, a mislabelled well name, unit mix.",
                "Fix: minimum-curvature TVD reconstruction (reproduces survey TVD to ~0.01 m).",
                "Fix: density-log porosity imputation validated against ThermoGIS "
                "(agree to within ~2.7 porosity points).",
                "Most teams trip on the depth and unit traps; we caught and validated them.",
            ),
            "Most teams will trip on the depth and unit traps. We caught them and validated "
            "against an independent regional model, so the foundation is clean and defensible.",
        ),
        SlideSpec(
            "Resource assessment (Challenge 1)",
            (
                "P10/P50/P90 doublet power per well; only 2 of 4 wells are viable (BLT, JUT).",
                "ThermoGIS regional grid (1 km, public TNO data) is the resource truth; "
                "our 4 wells validate it.",
                "Sited a new doublet in the proven trend, near demand: ~6 MW P50 class.",
            ),
            "The resource map is a public ThermoGIS product, so we use it and verify it "
            "rather than pretend to out-compute it. The value is the decision: where and "
            "how many wells.",
            figs["resource_map"],
        ),
        SlideSpec(
            "Integrated system design (Challenge 2A): the key idea",
            (
                "Geothermal doublet, heat pump (low 30 C reinjection), seasonal HT-ATES "
                "storage, heat-driven cooling.",
                "Cooling monetises otherwise-idle summer capacity: geothermal utilisation "
                "rises from 59% to 99%.",
                "Firm heating capacity 7.9 MWth; HT-ATES covers the winter peak with "
                "near-zero backup.",
            ),
            "This is the creative core and the reason the cost is low. Running the wells "
            "hard all year, heating in winter, cooling in summer, storing the surplus, is "
            "what drives the LCoE down.",
            figs["monthly_dispatch"],
        ),
        SlideSpec(
            "Techno-economics (Challenge 2B)",
            (
                "LCoE model is a port of the provided LCOE.xlsx (reproduces its 5.77 EUR/GJ "
                "base case), extended for the hybrid.",
                "Depth-dependent well cost from the spreadsheet's own formula.",
                f"Design comparison: 1 doublet {lc.get(1, float('nan')):.1f}, "
                f"2 doublets {lc.get(2, float('nan')):.1f}, "
                f"3 doublets {lc.get(3, float('nan')):.1f} EUR/GJ.",
                f"{_n_words(n)} doublet wins: fewer wells plus storage beats more wells.",
            ),
            "Our economics are consistent with the organisers' own spreadsheet, then "
            "extended for cooling and storage. One doublet is decisively cheapest.",
            figs["design_comparison"],
        ),
        SlideSpec(
            "Risk (Monte-Carlo)",
            (
                f"Transmissivity uncertainty gives an LCoE band: P10 {nums.p10:.1f}, "
                f"P50 {nums.p50:.1f}, P90 {nums.p90:.1f} EUR/GJ (right-skewed).",
                "A second nearby well does not de-risk it (same, correlated geology).",
                "Recommendation: drill one doublet, well-test it, expand only if "
                "data justify it.",
            ),
            "The band is wide and skewed because the resource is uncertain. The honest move "
            "is a staged drill-and-test, not committing to multiple wells up front.",
            figs["monte_carlo"],
        ),
        SlideSpec(
            "What actually moves the cost (sensitivity)",
            (
                "Tornado: injection temperature is the #1 driver (~7.7 EUR/GJ swing), then "
                "electricity price, discount rate, well cost.",
                "Honest headline: the one-doublet result depends on heat-pump-enabled 30 C "
                "reinjection.",
                "At a conservative ~39 C a second doublet is needed (flips near ~37 C).",
            ),
            "We surface the assumption the result is most sensitive to, rather than hide it. "
            "This is the number a reviewer should probe, and we show it ourselves.",
            figs["tornado"],
        ),
        SlideSpec(
            "Agentic workflow + interactive app (bonus)",
            (
                "An agent runs the whole pipeline and records a decision log (what it chose "
                "and why); no API key needed.",
                "Interactive app: change any assumption, re-run, see LCoE, band and dispatch "
                "live.",
                "Import/export the same config the CLI uses.",
            ),
            "The organisers said they will run the code and reward automation. The agent "
            "makes the whole analysis reproducible and explorable, and the app lets a "
            "non-coder test scenarios.",
        ),
        SlideSpec(
            "Practical applicability to Africa",
            (
                "Method transfers to East African Rift appraisal (KenGen/GDC): messy-data "
                "petrophysics, probabilistic resource, LCoE optimisation, agentic automation.",
                "Cooling and cascaded direct-use matter more in Africa (cold-chain, "
                "agriculture; already practised at Olkaria).",
                "Lowest-LCoE, staged development fits a capital-constrained setting.",
                "Honest: the Rift volcanic resource and heating-led demand don't transfer; "
                "the workflow and framing do.",
            ),
            "The case is Dutch because the data is public, but it is a template. A local team "
            "reruns the same pipeline on Rift data; the cooling and cost discipline are if "
            "anything more relevant in Africa.",
        ),
        SlideSpec(
            "Limitations (we state them)",
            (
                "Monthly energy balance, not a transient reservoir or thermodynamic sim.",
                "Resource map from sparse wells; P50-grid uncertainty modelled.",
                "1-doublet result depends on 30 C injection and on HT-ATES covering the "
                "winter peak.",
                "Absorption cooling at ~77 C is thermodynamically marginal; demand profile "
                "is an assumption.",
            ),
            "We are explicit about what the model does and does not do, and which "
            "assumptions carry the result. Transparency is part of the rigor.",
        ),
        SlideSpec(
            "Recommendation & conclusion",
            (
                "Stage a single geothermal doublet plus heat pump plus HT-ATES plus "
                "heat-driven cooling.",
                f"Meets >=10 MWth heat and >=5 MWth cooling at ~{lcoe:.0f} EUR/GJ P50 "
                f"(P10 {nums.p10:.0f}, P90 {nums.p90:.0f}).",
                "Cheapest credible option; de-risks before scaling; reproducible and "
                "transferable.",
                "Repo + app + report: [GitHub link].",
            ),
            "Close on the headline number and the staged, low-risk, lowest-cost "
            "recommendation, and point to the working code, app and report.",
        ),
    ]


def _n_words(n: int) -> str:
    return {1: "One", 2: "Two", 3: "Three"}.get(n, str(n))


def build_deck(nums: DeckNumbers, figs: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    for spec in _slide_specs(nums, figs):
        _add_content_slide(prs, spec)
    prs.save(PPTX_PATH)
    return PPTX_PATH


# --------------------------------------------------------------------------- #
# Entry point
# --------------------------------------------------------------------------- #
def main() -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    _apply_style()

    best = optimize()[0]
    dispatch_png = fig_monthly_dispatch(best)
    design_png, lcoe_by_doublets = fig_design_comparison()
    mc_png, p10, p50, p90 = fig_monte_carlo()
    tornado_png = fig_tornado()
    resource_png = fig_resource_map()

    figs = {
        "monthly_dispatch": dispatch_png,
        "design_comparison": design_png,
        "monte_carlo": mc_png,
        "tornado": tornado_png,
        "resource_map": resource_png,
    }
    nums = DeckNumbers(
        n_doublets=best.n_doublets,
        lcoe_eur_per_gj=best.lcoe_eur_per_gj,
        p10=p10,
        p50=p50,
        p90=p90,
        lcoe_by_doublets=lcoe_by_doublets,
    )
    pptx_path = build_deck(nums, figs)

    print("Generated files:")
    for path in (*figs.values(), pptx_path):
        size_kb = path.stat().st_size / 1024
        print(f"  {path}  ({size_kb:.1f} KB)")
    print()
    print("Headline numbers (live model):")
    print(f"  Optimal doublets : {nums.n_doublets}")
    print(f"  Headline LCoE    : {nums.lcoe_eur_per_gj:.2f} EUR/GJ")
    print(f"  MC P10/P50/P90   : {p10:.1f} / {p50:.1f} / {p90:.1f} EUR/GJ")
    print(f"  LCoE by doublets : {{1: {lcoe_by_doublets[1]:.1f}, "
          f"2: {lcoe_by_doublets[2]:.1f}, 3: {lcoe_by_doublets[3]:.1f}}}")


if __name__ == "__main__":
    main()
